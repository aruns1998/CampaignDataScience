from pptx import Presentation
from pptx.util import Inches
import os

os.makedirs('output', exist_ok=True)

prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Data-Driven Political Campaign"
slide.placeholders[1].text = "Thiruvidaimaruthur Constituency (2006–2024)"

# Add pie chart image
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Vote Share 2021"
slide.shapes.add_picture('output/vote_share_2021.png', Inches(1), Inches(1), width=Inches(6))

prs.save('output/campaign_analysis.pptx')
print("PPT generated!")
